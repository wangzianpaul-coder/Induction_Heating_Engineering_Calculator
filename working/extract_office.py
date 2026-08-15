from __future__ import annotations

import json
import os
from pathlib import Path
from zipfile import ZipFile

from docx import Document
from lxml import etree
from pptx import Presentation


def extract_docx(path: Path, output_dir: Path) -> dict:
    document = Document(str(path))
    records: list[dict] = []
    for index, paragraph in enumerate(document.paragraphs, 1):
        text = paragraph.text.strip()
        if text:
            records.append({"kind": "paragraph", "index": index, "text": text})
    for table_index, table in enumerate(document.tables, 1):
        rows = []
        for row_index, row in enumerate(table.rows, 1):
            rows.append(
                {
                    "row": row_index,
                    "cells": [cell.text.replace("\n", " / ") for cell in row.cells],
                }
            )
        records.append(
            {
                "kind": "table",
                "index": table_index,
                "rows": rows,
                "row_count": len(table.rows),
                "column_count": len(table.columns),
            }
        )

    media_dir = output_dir / "docx_images"
    media_dir.mkdir(parents=True, exist_ok=True)
    with ZipFile(path) as archive:
        media_names = [name for name in archive.namelist() if name.startswith("word/media/")]
        for name in media_names:
            (media_dir / Path(name).name).write_bytes(archive.read(name))
        xml = etree.fromstring(archive.read("word/document.xml"))
        namespaces = {
            "w": "http://schemas.openxmlformats.org/wordprocessingml/2006/main",
            "m": "http://schemas.openxmlformats.org/officeDocument/2006/math",
        }
        math_texts = []
        for math_node in xml.xpath("//m:oMath | //m:oMathPara", namespaces=namespaces):
            parts = math_node.xpath(".//m:t/text() | .//w:t/text()", namespaces=namespaces)
            joined = "".join(parts).strip()
            if joined and joined not in math_texts:
                math_texts.append(joined)
    result = {
        "file": str(path),
        "paragraph_count": len(document.paragraphs),
        "table_count": len(document.tables),
        "image_count": len(media_names),
        "records": records,
        "math_texts": math_texts,
    }
    (output_dir / "FURNACE_CONDITIONS_EXTRACT.json").write_text(
        json.dumps(result, ensure_ascii=False, indent=2), encoding="utf-8"
    )
    return result


def extract_pptx(path: Path, output_dir: Path) -> dict:
    presentation = Presentation(str(path))
    slides: list[dict] = []
    for slide_number, slide in enumerate(presentation.slides, 1):
        texts: list[str] = []
        tables: list[list[list[str]]] = []
        for shape in slide.shapes:
            text = getattr(shape, "text", "")
            if text and text.strip():
                texts.append(text.strip().replace("\x0b", "\n"))
            if getattr(shape, "has_table", False):
                table_rows = []
                for row in shape.table.rows:
                    table_rows.append([cell.text.replace("\n", " / ") for cell in row.cells])
                tables.append(table_rows)
        notes = ""
        try:
            notes = slide.notes_slide.notes_text_frame.text.strip()
        except Exception:
            pass
        slides.append(
            {
                "slide_number": slide_number,
                "texts": texts,
                "tables": tables,
                "notes": notes,
                "shape_count": len(slide.shapes),
            }
        )
    result = {"file": str(path), "slide_count": len(slides), "slides": slides}
    (output_dir / "PROJECT_PRESENTATION_EXTRACT.json").write_text(
        json.dumps(result, ensure_ascii=False, indent=2), encoding="utf-8"
    )
    return result


def main() -> None:
    output_dir = Path(os.environ["IH_OUTPUT_DIR"])
    output_dir.mkdir(parents=True, exist_ok=True)
    docx_result = extract_docx(Path(os.environ["IH_DOCX_PATH"]), output_dir)
    pptx_result = extract_pptx(Path(os.environ["IH_PPTX_PATH"]), output_dir)
    print(
        json.dumps(
            {
                "docx": {
                    "paragraph_count": docx_result["paragraph_count"],
                    "table_count": docx_result["table_count"],
                    "image_count": docx_result["image_count"],
                    "math_count": len(docx_result["math_texts"]),
                },
                "pptx": {"slide_count": pptx_result["slide_count"]},
            },
            ensure_ascii=False,
        )
    )


if __name__ == "__main__":
    main()
